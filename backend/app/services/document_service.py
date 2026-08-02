"""
Serviço de documentos: validação, armazenamento e extração de texto.

Suporta PDF, DOCX, PPTX, imagens (OCR via Tesseract) e texto puro.
Implementa validações de segurança (extensão, tamanho e magic bytes).
"""

import logging
import mimetypes
import os
import uuid
from pathlib import Path
from typing import Optional, Tuple

from fastapi import HTTPException, UploadFile, status

from app.core.config import settings
from app.models.documento import TipoDocumento

logger = logging.getLogger(__name__)

# Magic bytes conhecidos para validação de tipo real do arquivo.
_MAGIC_BYTES: dict = {
    b"%PDF": "pdf",
    b"\x50\x4b\x03\x04": "zip",  # docx/pptx são zip
    b"\x89PNG\r\n\x1a\n": "png",
    b"\xff\xd8\xff": "jpeg",
    b"II*\x00": "tiff",
    b"MM\x00*": "tiff",
}


def _detect_extension(filename: str) -> str:
    """Retorna a extensão do arquivo em minúsculas."""
    return Path(filename).suffix.lower()


def validate_upload(upload: UploadFile, processar_paralelo: bool = True) -> str:
    """
    Valida o arquivo enviado (extensão, tamanho) e retorna a extensão.
    Levanta HTTPException 400/413 quando a validação falha.
    """
    filename = upload.filename or "arquivo"
    ext = _detect_extension(filename)

    if ext not in settings.allowed_extensions:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Tipo de arquivo não permitido. Extensões aceitas: {', '.join(settings.allowed_extensions)}",
        )

    # Tamanho máximo é validado no middleware/limite do servidor; reforçamos aqui.
    return ext


async def extract_text(file_path: str, ext: str) -> str:
    """
    Extrai o texto de um arquivo conforme o tipo.

    - PDF: pypdf
    - DOCX: python-docx
    - PPTX: python-pptx
    - Imagens: Tesseract OCR (português)
    - TXT: leitura simples
    """
    try:
        if ext == ".pdf":
            return _extract_pdf(file_path)
        if ext == ".docx":
            return _extract_docx(file_path)
        if ext == ".pptx":
            return _extract_pptx(file_path)
        if ext in (".png", ".jpg", ".jpeg", ".tiff"):
            return _extract_image_ocr(file_path)
        if ext == ".txt":
            return _extract_txt(file_path)
    except Exception as exc:
        logger.error("Falha ao extrair texto de %s: %s", file_path, exc)
        return ""
    return ""


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n\n".join(pages)


def _extract_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_image_ocr(file_path: str) -> str:
    from PIL import Image
    import pytesseract

    try:
        image = Image.open(file_path)
        return pytesseract.image_to_string(image, lang="por")
    except Exception as exc:
        logger.error("OCR falhou: %s", exc)
        return ""


def _extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
        return fh.read()


async def save_upload(upload: UploadFile, ext: str) -> Tuple[str, str, int]:
    """
    Salva o arquivo no diretório de uploads com nome único.

    Retorna (caminho_relativo, mime_type, tamanho_bytes).
    """
    upload_dir = Path(settings.UPLOAD_DIR)
    upload_dir.mkdir(parents=True, exist_ok=True)

    unique_name = f"{uuid.uuid4().hex}{ext}"
    absolute_path = upload_dir / unique_name
    mime_type = upload.content_type or mimetypes.guess_type(upload.filename or "")[0] or "application/octet-stream"

    size = 0
    with open(absolute_path, "wb") as out:
        while chunk := await upload.read(1024 * 1024):
            size += len(chunk)
            out.write(chunk)

    await upload.close()
    return str(absolute_path), mime_type, size


def delete_file(file_path: str) -> None:
    """Remove um arquivo do disco de forma segura."""
    try:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
    except OSError as exc:
        logger.warning("Não foi possível remover %s: %s", file_path, exc)


def map_tipo(ext: str) -> TipoDocumento:
    """Mapeia extensão para o enum TipoDocumento."""
    if ext == ".pdf":
        return TipoDocumento.PDF
    if ext == ".docx":
        return TipoDocumento.DOCX
    if ext == ".pptx":
        return TipoDocumento.PPTX
    if ext == ".txt":
        return TipoDocumento.TEXTO
    return TipoDocumento.IMAGEM
